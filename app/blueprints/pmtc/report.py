"""
report.py -- PPT generation skeleton.
Before coding: run the shape audit from PPT_CONVENTIONS.md Part 2.
Identify pre-filled shapes (never push to these) vs empty shapes.
"""
import matplotlib
matplotlib.use('Agg')  # must be at module level -- before any pyplot import

import io, shutil, tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

TEMPLATE = Path(__file__).parent.parent.parent.parent / 'TEMPLATE.pptx'

def generate_report(kpis, profile, priorities, investment=None):
    """
    Populate the PPT template with user data and return bytes.
    Steps:
    1. Identify active benefits from priorities
    2. Slide deletion pre-pass (collect, then delete in reverse)
    3. Populate empty shapes by name
    4. Insert matplotlib charts into img_ placeholders
    """
    prs = Presentation(TEMPLATE)
    active_benefits = _active_benefits(priorities)

    # --- 1. Slide deletion pre-pass ---
    # Collect indices of slides to delete (inactive benefit markers)
    to_delete = []
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.name.startswith('marker_'):
                benefit_id = shape.name[len('marker_'):]
                if benefit_id not in active_benefits:
                    to_delete.append(idx)
    # Delete in reverse order to preserve indices
    for idx in sorted(set(to_delete), reverse=True):
        rId = prs.slides._sldIdLst[idx].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[idx]

    # --- 2. Populate shapes ---
    company = profile.get('company', '[Company]')
    for slide in prs.slides:
        for shape in slide.shapes:
            _push(shape, 'txt_companyName', company)
            _push(shape, 'txt_ROI', f"{kpis.get('roi', 0):.0f}x")
            # TODO: add remaining shape pushes from actual template audit

    # --- 3. Insert charts ---
    # TODO: generate and insert matplotlib charts into img_ shapes

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

def _push(shape, name, value):
    """Push text to a named shape if it matches and has a text frame."""
    if shape.name == name and shape.has_text_frame:
        tf = shape.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = str(value)
        else:
            tf.paragraphs[0].text = str(value)

def _active_benefits(priorities):
    """Return set of benefit IDs active based on priorities (High or Medium)."""
    # TODO: implement from challenge-benefit matrix in CLAUDE.md
    # Example:
    # from app.blueprints.pmtc.calculator import BENEFIT_MATRIX
    # active = set()
    # for ch_id, level in priorities.items():
    #     if level in ('High', 'Medium'):
    #         for b in BENEFIT_MATRIX.get(ch_id, []):
    #             active.add(b)
    # return active
    return set()
