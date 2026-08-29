"""Build the deck that becomes the PDF.

**This is the placeholder.** Ben's real PMTC report template does not exist
yet (PROJECT_STATE.md Open Item #2), so rather than block the whole mail
pipeline on it, this draws a single slide from scratch with python-pptx and
fills it with the visitor's real computed results. Everything around it -- the
container, LibreOffice, the fonts, the raw-MIME attachment, SES, the async
invoke from Flask -- is the final architecture and is exercised for real.

The point of drawing something real rather than sending an empty PDF is that
the parts most likely to break on the day the template lands are the parts a
placeholder still exercises: that python-pptx output converts at all, how long
the conversion takes on this memory setting, whether the fonts resolve, and
whether a multi-megabyte attachment survives SES.

See README.md, "When the real deck arrives", for what replaces this. The
`generate(data, out_path)` signature is deliberately the one the real
generator will keep, so `handler.py` does not change when it does.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# The tool's own palette, lifted from results.html's custom properties so the
# placeholder is recognisably the same product rather than default Office
# blue. `ink` is the body colour, `signal` is the highlight behind the current
# breadcrumb step, `navy_deep` is the results band.
INK = RGBColor(0x1E, 0x1E, 0x1E)
INK_SOFT = RGBColor(0x3A, 0x3A, 0x36)
SLATE = RGBColor(0x6E, 0x6E, 0x68)
SIGNAL = RGBColor(0xDF, 0xDA, 0x7A)
NAVY_DEEP = RGBColor(0x06, 0x16, 0x27)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9, which is what any deck K1x supplies will almost certainly be.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Named here rather than inline so the substitution rule in fonts.conf and the
# family this asks for cannot drift apart.
FONT = "Outfit"


def _textbox(slide, left, top, width, height, text, *, size, bold=False,
             color=INK, align=PP_ALIGN.LEFT, spacing=None):
    """One text box, with the paragraph properties set on every run.

    python-pptx applies font settings per run, not per shape, so a helper that
    forgets a run leaves that line in the theme default -- which in a
    from-scratch presentation is Calibri, not the family fonts.conf has a rule
    for. Everything drawn here goes through this function for that reason.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    # Zero the insets. python-pptx defaults a text box to a 0.1in left and
    # right margin, so a shape placed at 0.75in does not start its text at
    # 0.75in -- and two boxes at different x values that were meant to share a
    # left edge end up visibly out of line. Zeroing makes the geometry here
    # mean what it says.
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    lines = text.split("\n")
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        if spacing is not None:
            para.line_spacing = spacing
        run = para.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _band(slide, left, top, width, height, color):
    """A flat colour rectangle with no outline, used as a rule or a fill."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def generate(data, out_path):
    """Write the deck for one assessment to `out_path`.

    `data` is the payload `emailer.py` sends: the whole `run_calculation()`
    result under "results", plus the lead-capture fields under "lead". Read
    defensively -- a missing key here must not be the reason a visitor never
    gets their report.
    """
    results = data.get("results") or {}
    lead = data.get("lead") or {}

    company = results.get("company") or lead.get("company") or "Your firm"
    industry = results.get("industry") or ""
    your_score = results.get("your_score")
    peer_score = results.get("peer_score")
    peer_count = results.get("peer_count")
    band_name = results.get("band_name") or ""
    band_subtitle = results.get("band_subtitle") or ""
    narrative = results.get("narrative") or ""
    strengths = results.get("strengths") or []
    gaps = results.get("gaps") or []

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    # Layout 6 is the blank one in the default template. Anything else brings
    # placeholder shapes that would need deleting.
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _band(slide, 0, 0, SLIDE_W, Inches(1.45), NAVY_DEEP)
    _textbox(slide, Inches(0.75), Inches(0.34), Inches(11.8), Inches(0.4),
             "K1X PRIVATE MARKET TAX CAPABILITY ASSESSMENT",
             size=12, bold=True, color=SIGNAL)
    _textbox(slide, Inches(0.75), Inches(0.66), Inches(11.8), Inches(0.6),
             company, size=28, bold=True, color=PAPER)

    top = Inches(1.95)
    if industry:
        _textbox(slide, Inches(0.75), top, Inches(11.8), Inches(0.3),
                 industry.upper(), size=11, bold=True, color=SLATE)

    # The score, the peer score and the band. The three numbers the Results
    # page leads with, in the same order.
    _textbox(slide, Inches(0.75), Inches(2.35), Inches(3.0), Inches(1.2),
             "" if your_score is None else "{:.1f}".format(your_score),
             size=72, bold=True, color=INK)
    _textbox(slide, Inches(0.75), Inches(3.55), Inches(3.0), Inches(0.3),
             "YOUR MATURITY SCORE (0-5)", size=10, bold=True, color=SLATE)

    peer_line = "Peer leaders: {}".format(
        "" if peer_score is None else "{:.1f}".format(peer_score))
    if peer_count:
        peer_line += "   (n={})".format(peer_count)
    _textbox(slide, Inches(0.75), Inches(3.95), Inches(4.0), Inches(0.3),
             peer_line, size=13, color=INK_SOFT)

    _textbox(slide, Inches(4.6), Inches(2.4), Inches(8.0), Inches(0.5),
             band_name, size=26, bold=True, color=INK)
    if band_subtitle:
        _textbox(slide, Inches(4.6), Inches(2.95), Inches(8.0), Inches(0.4),
                 band_subtitle, size=15, color=SLATE)
    _textbox(slide, Inches(4.6), Inches(3.45), Inches(8.0), Inches(1.6),
             narrative, size=12, color=INK_SOFT, spacing=1.3)

    _band(slide, Inches(0.75), Inches(5.25), Inches(11.83), Emu(12700), SIGNAL)

    _textbox(slide, Inches(0.75), Inches(5.45), Inches(5.5), Inches(0.3),
             "WHERE YOU SCORED HIGHEST", size=11, bold=True, color=SLATE)
    strength_lines = "\n".join(
        "{}   {}".format(
            item.get("name", ""),
            "" if item.get("score") is None else "{:.1f}".format(item["score"]),
        )
        for item in strengths
    )
    _textbox(slide, Inches(0.75), Inches(5.8), Inches(5.5), Inches(1.2),
             strength_lines, size=13, color=INK, spacing=1.35)

    _textbox(slide, Inches(6.9), Inches(5.45), Inches(5.7), Inches(0.3),
             "PRIORITY AREAS TO IMPROVE", size=11, bold=True, color=SLATE)
    gap_lines = "\n".join(
        "{}   {}".format(
            item.get("name", ""),
            "" if item.get("delta") is None else "{:+.1f} vs peers".format(item["delta"]),
        )
        for item in gaps
    )
    _textbox(slide, Inches(6.9), Inches(5.8), Inches(5.7), Inches(1.2),
             gap_lines, size=13, color=INK, spacing=1.35)

    # Said on the artifact itself, not only in the covering email. A PDF gets
    # forwarded and read away from the message that carried it, and a reader
    # who cannot tell a placeholder from the deliverable will assume this is
    # the deliverable.
    _textbox(slide, Inches(0.75), Inches(7.0), Inches(11.8), Inches(0.3),
             "Placeholder layout. The full K1x report template is in "
             "preparation; the figures above are your real results.",
             size=9, color=SLATE)

    prs.save(out_path)
    return out_path
